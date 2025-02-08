import streamlit as st
import requests
from PyPDF2 import PdfReader
from docx import Document
from pptx import Presentation
import json
import io
import base64 #Converts file content to a downloadable format.

# Configuration
API_URL = "http://localhost:3000"
CHATFLOW_ID = "498be727-db09-4e73-858e-9eaab7998ec9"
WHATSAPP_TEST_NUMBER = "+15551756586"

def upload_content_to_flowise(file_name, file_content):
    try:
        mime_type = (
            "application/pdf" if file_name.endswith(".pdf") else
            "application/vnd.openxmlformats-officedocument.wordprocessingml.document" if file_name.endswith(".docx") else
            "application/vnd.openxmlformats-officedocument.presentationml.presentation" if file_name.endswith(".pptx") else
            "text/plain"
        )
        
        form_data = {"chatId": "upload-session-" + file_name}
        files = {"files": (file_name, file_content, mime_type)}

        url = API_URL + "/api/v1/vector/upsert/" + CHATFLOW_ID
        response = requests.post(url, files=files, data=form_data)
        response_json = response.json()

        if response.status_code == 200:
            if response_json.get("numAdded", 0) > 0:
                return {"success": True, "message": "Successfully uploaded " + str(response_json["numAdded"]) + " chunks"}
            else:
                if "error" in response_json:
                    return {"error": "Server error: " + response_json["error"]}
                elif "message" in response_json:
                    if "File processed but no vectors were added" in response_json["message"]:
                        return {"duplicate": True}
                    return {"error": "Server message: " + response_json["message"]}
                else:
                    return {"duplicate": True}
        else:
            return {"error": "Upload failed with status code: " + str(response.status_code)}

    except Exception as e:
        return {"error": str(e)}

def process_file(uploaded_file):
    """
    Process uploaded file and return its extracted text and content.
    """
    try:
        file_name = uploaded_file.name.lower()
        file_content = uploaded_file.read()
        text_content = ""

        # Process PDF files
        if file_name.endswith(".pdf"):
            reader = PdfReader(io.BytesIO(file_content))
            text_content = "\n".join(page.extract_text() for page in reader.pages if page.extract_text())

        # Process DOCX files
        elif file_name.endswith(".docx"):
            doc = Document(io.BytesIO(file_content))
            text_content = "\n".join(paragraph.text for paragraph in doc.paragraphs)

        # Process PPTX files
        elif file_name.endswith(".pptx"):
            prs = Presentation(io.BytesIO(file_content))
            text_content = "\n".join(
                shape.text 
                for slide in prs.slides 
                for shape in slide.shapes 
                if hasattr(shape, "text")
            )

        # Ensure text is extracted
        if not text_content.strip():
            raise ValueError("The file appears to be blank.")

        return file_content, text_content

    except Exception as e:
        raise ValueError(str(e))

# Streamlit UI 
st.title("Gen AI Mental Health Support ChatBot 🤖")
st.write("Welcome! This chatbot provides preliminary support and resources related to mental health.")

# Initialize session state for storing uploaded files
if "uploaded_files_list" not in st.session_state:
    st.session_state["uploaded_files_list"] = []
if "uploaded_files_content" not in st.session_state:
    st.session_state["uploaded_files_content"] = {}  # Store content as key-value (name: content)

# File Upload Section
st.subheader("Upload Mental Health Resources (PDF, Word, or PPTX)")
uploaded_files = st.file_uploader("Choose files to upload", type=["pdf", "docx", "pptx", "txt"], accept_multiple_files=True)

if st.button("Upload Files"):
    if not uploaded_files:
        st.error("Please upload at least one file.")
    else:
        progress_bar = st.progress(0)
        total_files = len(uploaded_files)

        for index, uploaded_file in enumerate(uploaded_files):
            try:
                file_name = uploaded_file.name
                base_name = ".".join(file_name.split(".")[:-1])  # Remove extension
                
                if any(base_name in existing for existing in st.session_state["uploaded_files_list"]):
                    st.warning("'" + file_name + "' - Already Uploaded")
                    continue 

                # Process the file and get content and extracted text
                file_content, extracted_text = process_file(uploaded_file)

                # Upload to Flowise
                with st.spinner("Processing " + uploaded_file.name + "..."):
                    upload_response = upload_content_to_flowise(uploaded_file.name, file_content)

                    if "duplicate" in upload_response:
                        st.warning("'" + uploaded_file.name + "' - Already Uploaded")
                    elif "error" in upload_response:
                        st.error("Failed to upload " + uploaded_file.name + ": " + upload_response["error"])
                    else:
                        st.success(uploaded_file.name + " uploaded successfully! " + upload_response.get("message", ""))

                        # Store file name and content in session state
                        if uploaded_file.name not in st.session_state["uploaded_files_list"]:
                            st.session_state["uploaded_files_list"].append(uploaded_file.name)
                            st.session_state["uploaded_files_content"][uploaded_file.name] = file_content  # Store content

                        # Display Extracted Text in a Collapsible Section
                        with st.expander("Extracted Content from " + uploaded_file.name):
                            st.text_area("Extracted Text", extracted_text, height=200)

            except Exception as e:
                st.error("Error processing " + uploaded_file.name + ": " + str(e))

            finally:
                # Update progress bar
                progress_bar.progress((index + 1) / total_files)

# Sidebar Display
st.sidebar.markdown("📂 Uploaded Files")
st.sidebar.markdown("<hr style='border: 1px solid #fff;'>", unsafe_allow_html=True)

if st.session_state["uploaded_files_list"]:
    for file_name in st.session_state["uploaded_files_list"]:
        file_content = st.session_state["uploaded_files_content"].get(file_name, None)  # Retrieve stored content

        if file_content:
            # Convert file content to base64 for download link
            b64_file = base64.b64encode(file_content).decode()
            file_extension = file_name.split(".")[-1]

            # File type icons
            file_icon = {
                "pdf": "📄",
                "docx": "📝",
                "pptx": "📊",
            }.get(file_extension, "📁")

            # Generate clickable link
            href = f'<a href="data:application/octet-stream;base64,{b64_file}" download="{file_name}" style="color:white; text-decoration: none; font-weight: bold; display: flex; align-items: center; gap: 8px; padding: 6px 10px; border-radius: 8px; background: rgba(255,255,255,0.2); margin-bottom: 6px; transition: 0.3s;" onmouseover="this.style.background=\'rgba(255,255,255,0.4)\'" onmouseout="this.style.background=\'rgba(255,255,255,0.2)\'">{file_icon} {file_name}</a>'

            st.sidebar.markdown(href, unsafe_allow_html=True)
else:
    st.sidebar.write("🔹 No files uploaded yet.")

st.markdown(
    """
    <style>
    /* Sidebar Background - Soft Blue Gradient */
    [data-testid="stSidebar"] {
        background: linear-gradient(to bottom, #0078D7, #00BFFF);
        color: white;
        padding: 20px;
    }
    
    /* Sidebar Text and Headers */
    [data-testid="stSidebar"] h3, [data-testid="stSidebar"] h4, [data-testid="stSidebar"] p {
        color: white !important;
        font-weight: bold;
    }

    /* Uploaded Files List - Rounded Box with Hover */
    [data-testid="stSidebar"] a {
        display: block;
        padding: 8px 12px;
        margin: 6px 0;
        border-radius: 8px;
        background: rgba(255, 255, 255, 0.2);
        color: white;
        text-decoration: none;
        font-weight: bold;
        transition: 0.3s;
    }

    [data-testid="stSidebar"] a:hover {
        background: rgba(255, 255, 255, 0.4);
    }
    
    /* Sidebar Folder Icon */
    [data-testid="stSidebar"] .icon-folder {
        font-size: 22px;
        margin-right: 8px;
    }
    </style>
    """,
    unsafe_allow_html=True
)



st.markdown(
    """
    <style>
    /* Main Page Background */
    html, body, .stApp {
        background: linear-gradient(to bottom, #ff758c, #ff7eb3) !important;
        height: 100vh !important;
        width: 100vw !important;
        margin: 0 !important;
        padding: 0 !important;
        overflow-x: hidden !important;
        display: flex !important;
        justify-content: center !important;
        align-items: center !important;
    }
    /* Ensure Content is Centered */
    .block-container {
        max-width: 700px !important;
        text-align: center !important;
        padding-top: 5rem !important;
        margin: auto !important;
    }
    /* File Uploader Section */
    [data-testid="stFileUploader"] {
        background: rgba(255, 255, 255, 0.95) !important;
        border-radius: 12px !important;
        padding: 15px !important;
        width: 90% !important;
        max-width: 500px !important;
        margin: auto !important;
        box-shadow: 2px 2px 10px rgba(0,0,0,0.1) !important;
    }
    /* Uploaded Files Box */
    .stAlert {
        font-size: 16px !important;
        font-weight: bold !important;
        color: #333333 !important; /* Dark text for contrast */
        background: rgba(255, 255, 255, 0.9) !important; /* Light background */
        border-radius: 12px !important;
        padding: 15px !important;
        width: 80% !important;
        max-width: 600px !important;
        margin: 10px auto !important; /* Adds spacing */
        box-shadow: 2px 2px 10px rgba(0,0,0,0.1) !important; /* Subtle shadow */
    }
    /* Progress Bar Styling */
    div[data-testid="stProgress"] {
        width: 80% !important;
        margin: auto !important;
    }
    /* Button Styling */
    .stButton>button {
        width: 80% !important;
        max-width: 300px !important;
        font-size: 16px;
        padding: 12px;
        border-radius: 8px;
        cursor: pointer;
        transition: transform 0.2s ease-in-out;
        background: linear-gradient(to right, #ff4081, #ff80ab) !important;
        color: white !important;
        border: none !important;
    }
    .stButton>button:hover {
        transform: scale(1.05);
        background: linear-gradient(to right, #f50057, #ff80ab) !important;
    }
    /* Keep the Streamlit header visible */
    header {
        visibility: visible !important;
        height: auto !important;
    }
    </style>
    """,
    unsafe_allow_html=True
)
# WhatsApp Floating Button
whatsapp_url = f"https://wa.me/{WHATSAPP_TEST_NUMBER}"
st.markdown(
    f"""
    <style>
        .whatsapp-button {{
            position: fixed;
            bottom: 20px;
            right: 20px;
            background-color: #25D366;
            color: white;
            border: none;
            padding: 15px 25px;
            border-radius: 50px;
            font-size: 16px;
            cursor: pointer;
            text-decoration: none;
            box-shadow: 2px 2px 5px rgba(0,0,0,0.3);
        }}
    </style>
    <a class="whatsapp-button" href="{whatsapp_url}" target="_blank">💬 Chat on WhatsApp</a>
    """,
    unsafe_allow_html=True
)





